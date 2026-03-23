"""
Core catalog generation logic.

Produces a product catalog .pptx from an Excel DataFrame, a slide template,
and an images directory. Designed to be called from app.py (Streamlit) or CLI.
"""

import os
import copy
import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Config ───────────────────────────────────────────────────────────────────
TEMPLATE_SLIDE_IDX = 1  # 0-based; slide 2

# Main image
TEXT_CENTER = Inches(4)
MAIN_IMG_VCENTER = Inches(2.886)
MAIN_IMG_SIZE = 5.1

# Close-up circle layouts: list of (left, top) positions per layout
CIRCLE_LAYOUT_2 = {
    "size": Inches(2.7),
    "positions": [
        (Inches(8.40), Inches(0.77)),
        (Inches(8.40), Inches(4.13)),
    ],
}
CIRCLE_LAYOUT_3 = {
    "size": Inches(2.45),
    "positions": [
        (Inches(6.542), Inches(1.108)),
        (Inches(9.764), Inches(1.155)),
        (Inches(8.164), Inches(4.098)),
    ],
}
CIRCLE_LAYOUT_4 = {
    "size": Inches(2.2),
    "positions": [
        (Inches(7.064), Inches(0.736)),
        (Inches(9.930), Inches(0.736)),
        (Inches(7.064), Inches(3.750)),
        (Inches(9.930), Inches(3.730)),
    ],
}

# Table config
TABLE_WIDTH_SMALL = Inches(4.80)
TABLE_WIDTH_LARGE = Inches(6.36)
TABLE_TOP = Inches(6.15)
TABLE_ROW_HEIGHT = Inches(0.40)
TABLE_FONT_NAME = "Georgia Pro Light"
TABLE_FONT_SIZE = Pt(12)
TABLE_HEADER_FILL = RGBColor(0x3B, 0x45, 0x40)
TABLE_WIDTH_BOOST = 0.2


# ── Helpers ──────────────────────────────────────────────────────────────────

def read_excel(path):
    df = pd.read_excel(path)
    df.columns = [c.strip().lower() for c in df.columns]
    return df


def find_image(folder, pattern):
    """Find first file in folder whose name contains pattern (case-insensitive).
    Prioritizes .png files over other formats."""
    if not os.path.isdir(folder):
        return None
    matches = [
        fname for fname in sorted(os.listdir(folder))
        if pattern.lower() in fname.lower()
        and os.path.isfile(os.path.join(folder, fname))
    ]
    if not matches:
        return None
    for fname in matches:
        if fname.lower().endswith(".png"):
            return os.path.join(folder, fname)
    return os.path.join(folder, matches[0])


def get_image_dimensions(image_path, target_inches=MAIN_IMG_SIZE):
    """Return (width, height) in Inches, fixing the dominant dimension to target_inches."""
    with PILImage.open(image_path) as img:
        w, h = img.size
    if h / w > 1:
        return Inches(target_inches * w / h), Inches(target_inches)
    else:
        return Inches(target_inches), Inches(target_inches * h / w)


def duplicate_slide(prs, src_idx):
    """Duplicate a slide by copying XML elements and image relationships."""
    src_slide = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src_slide.slide_layout)

    for ph in list(new_slide.placeholders):
        ph._element.getparent().remove(ph._element)

    for el in src_slide.shapes._spTree:
        tag_local = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        if tag_local in ("sp", "pic", "grpSp", "cxnSp"):
            new_slide.shapes._spTree.append(copy.deepcopy(el))

    for rel in src_slide.part.rels.values():
        if "image" in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    return new_slide


def update_text_of_textbox(slide, text_box_id, new_text):
    """Update the Nth text-bearing shape on a slide (1-indexed), preserving font formatting."""
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            if count == text_box_id:
                tf = shape.text_frame
                run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
                fmt = {
                    "name": run.font.name,
                    "size": run.font.size,
                    "bold": run.font.bold,
                    "color": run.font.color.rgb if run.font.color and run.font.color.rgb else None,
                }
                tf.clear()
                new_run = tf.paragraphs[0].add_run()
                new_run.text = new_text
                new_run.font.name = fmt["name"]
                new_run.font.size = fmt["size"]
                new_run.font.bold = fmt["bold"]
                if fmt["color"]:
                    new_run.font.color.rgb = fmt["color"]
                return


# ── Table helpers ────────────────────────────────────────────────────────────

def get_col(row, name):
    """Get a column value, returning None if missing or NaN."""
    val = row.get(name)
    if val is None or (isinstance(val, float) and pd.isna(val)):
        return None
    return val


def has_any_col(row, col_names):
    return any(get_col(row, c) is not None for c in col_names)


def format_dim_value(row):
    dims = []
    for col in ["length", "width", "height"]:
        val = get_col(row, col)
        if val is not None:
            dims.append(str(int(val)) if val == int(val) else str(val))
    return "x".join(dims) if dims else ""


def format_packing_value(row):
    parts = []
    inner = get_col(row, "inner carton")
    outer = get_col(row, "outer carton")
    unit = get_col(row, "unit")
    cbm = get_col(row, "cbm")

    qty_parts = []
    if inner is not None:
        qty_parts.append(str(int(inner)) if inner == int(inner) else str(inner))
    if outer is not None:
        qty_parts.append(str(int(outer)) if outer == int(outer) else str(outer))
    qty_str = "/".join(qty_parts)
    if unit is not None:
        qty_str = f"{qty_str} {unit}" if qty_str else str(unit)
    if qty_str:
        parts.append(qty_str)

    if cbm is not None:
        parts.append(f"{cbm} cbm")

    return " / ".join(parts)


def build_table_data(row, df_columns):
    DIM_COLS = ["length", "width", "height"]
    PACK_COLS = ["inner carton", "outer carton", "unit", "cbm"]
    columns = []

    if has_any_col(row, DIM_COLS):
        columns.append(("Dim. (cm)", format_dim_value(row)))

    price = get_col(row, "price")
    if price is not None:
        columns.append(("$", str(price)))

    if has_any_col(row, PACK_COLS):
        columns.append(("Packing", format_packing_value(row)))

    fob = get_col(row, "fob port")
    if fob is not None:
        columns.append(("FOB", str(fob)))

    return columns


def style_table_cell(cell, is_header=False):
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = TABLE_FONT_NAME
            run.font.size = TABLE_FONT_SIZE
    if is_header:
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_FILL


def add_product_table(slide, row, df_columns):
    table_data = build_table_data(row, df_columns)
    if not table_data:
        return

    num_cols = len(table_data)
    table_width = TABLE_WIDTH_SMALL if num_cols <= 3 else TABLE_WIDTH_LARGE
    table_left = TEXT_CENTER - table_width // 2

    table_shape = slide.shapes.add_table(2, num_cols, table_left, TABLE_TOP,
                                         table_width, TABLE_ROW_HEIGHT * 2)
    table = table_shape.table

    headers = [h for h, _ in table_data]
    col_width = int(table_width / num_cols)
    packing_bonus = int(col_width * TABLE_WIDTH_BOOST) if "Packing" in headers else 0
    non_packing_reduction = int(packing_bonus / max(num_cols - 1, 1)) if packing_bonus else 0
    for i in range(num_cols):
        if headers[i] == "Packing":
            table.columns[i].width = col_width + packing_bonus
        else:
            table.columns[i].width = col_width - non_packing_reduction

    for c, (header, value) in enumerate(table_data):
        hcell = table.cell(0, c)
        hcell.text = header
        style_table_cell(hcell, is_header=True)

        dcell = table.cell(1, c)
        dcell.text = value
        style_table_cell(dcell, is_header=False)


# ── Core ─────────────────────────────────────────────────────────────────────

def create_slide_from_row(prs, row, template_idx, df_columns, images_dir):
    """Create a new product slide from a DataFrame row."""
    product_id = str(get_col(row, "product_id") or "")
    img_folder = os.path.join(images_dir, product_id)

    has_folder = os.path.isdir(img_folder)
    main_img = find_image(img_folder, "IMG") if has_folder else None
    circles = []
    for n in range(1, 5):
        img = find_image(img_folder, f"circle{n}") if has_folder else None
        if img:
            circles.append(img)

    slide = duplicate_slide(prs, template_idx)

    update_text_of_textbox(slide, 1, product_id)

    # Remove the template table (duplicated from template slide)
    for shape in list(slide.shapes):
        if shape.has_table:
            shape._element.getparent().remove(shape._element)

    add_product_table(slide, row, df_columns)

    # Main image
    if main_img:
        img_w, img_h = get_image_dimensions(main_img)
        img_left = TEXT_CENTER - img_w // 2
        img_top = MAIN_IMG_VCENTER - img_h // 2
        slide.shapes.add_picture(main_img, img_left, img_top, img_w, img_h)

    # Circle images — select layout based on count
    if circles:
        layout = {2: CIRCLE_LAYOUT_2, 3: CIRCLE_LAYOUT_3, 4: CIRCLE_LAYOUT_4}.get(
            len(circles), CIRCLE_LAYOUT_2
        )
        size = layout["size"]
        for i, img_path in enumerate(circles):
            if i < len(layout["positions"]):
                left, top = layout["positions"][i]
                slide.shapes.add_picture(img_path, left, top, size, size)

    return has_folder


def generate_catalog(excel_path, template_path, images_dir, output_path):
    """
    Main entry point: read Excel → generate slides → save .pptx.
    Returns a summary dict for the UI.
    """
    df = read_excel(excel_path)
    prs = Presentation(template_path)
    df_columns = set(df.columns)

    missing_folders = []
    for _, row in df.iterrows():
        pid = str(row.get("product_id", ""))
        if not os.path.isdir(os.path.join(images_dir, pid)):
            missing_folders.append(pid)
        create_slide_from_row(prs, row, TEMPLATE_SLIDE_IDX, df_columns, images_dir)

    # Remove original template slide
    sldIdLst = prs._element.find(
        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst"
    )
    sldIdLst.remove(list(sldIdLst)[TEMPLATE_SLIDE_IDX])

    prs.save(output_path)

    total = len(df)
    found = total - len(missing_folders)
    return {
        "total_rows": total,
        "slides_created": total,
        "images_found": found,
        "images_pct": round(found / total * 100) if total else 0,
        "missing_folders": missing_folders,
        "output_path": output_path,
    }
